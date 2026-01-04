from pptx import Presentation
import os

def extract_text_from_pptx(filepath):
    """
    Extracts text from a PowerPoint file.
    """
    if not os.path.exists(filepath):
        print(f"File not found: {filepath}")
        return

    try:
        prs = Presentation(filepath)
        print(f"--- Analysis of: {os.path.basename(filepath)} ---")
        
        for i, slide in enumerate(prs.slides):
            print(f"\n[Slide {i+1}]")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(f"- {shape.text}")
                
                # Also check for text in tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text_frame.text for cell in row.cells])
                        print(f"  [Table Row] {row_text}")
                        
    except Exception as e:
        print(f"Error reading PPTX: {e}")

if __name__ == "__main__":
    pptx_path = r"d:\Projects\Privacy-Safe Cross-Company Data Insights\Prototype Submission Deck _ AI for Good Hackathon.pptx"
    extract_text_from_pptx(pptx_path)
