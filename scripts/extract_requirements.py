import os
from pptx import Presentation

def extract_text_from_pptx(file_path):
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        return

    print(f"--- Extracting text from: {os.path.basename(file_path)} ---")
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error opening PPTX: {e}")
        return

    full_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_title = "No Title"
        if slide.shapes.title:
            slide_title = slide.shapes.title.text
        
        slide_text.append(f"\n[Slide {i+1}: {slide_title}]")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                # Avoid duplicating the title if possible, but simple extraction is fine
                if shape.text != slide_title:
                    slide_text.append(shape.text)
        
        full_text.append("\n".join(slide_text))
    
    print("\n".join(full_text))

if __name__ == "__main__":
    pptx_path = r"d:\Projects\Privacy-Safe Cross-Company Data Insights\Prototype Submission Deck _ AI for Good Hackathon.pptx"
    extract_text_from_pptx(pptx_path)
